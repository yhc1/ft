"""
Generate PPTX slides for AIL Readiness Assessment.
One slide per dimension, formatted like the FT Strategies layout:
  - Title at top
  - Numbered prerequisite cards with circle badges
  - Use case table: Use Case | Data Inputs | Modelling Processes | Data Outputs

Outputs:
  AIL_Readiness_Assessment_EN.pptx
  AIL_Readiness_Assessment_ZH.pptx
"""

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree
from pptx.util import Cm, Pt

DATA_DIR = Path(__file__).parent / "web" / "data"

# ── Colors ───────────────────────────────────────────────────────────────────
BLACK = RGBColor(0x0D, 0x0D, 0x0D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x99, 0x99, 0x99)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)

# ── Slide dimensions (16:9 widescreen) ──────────────────────────────────────
SLIDE_W = 33.87   # cm
SLIDE_H = 19.05   # cm

# ── Layout grid ──────────────────────────────────────────────────────────────
ML = 1.3          # left margin
MR = 1.3          # right margin
CW = SLIDE_W - ML - MR   # content width = 31.27 cm

TITLE_TOP  = 0.4
TITLE_H    = 1.8

LABEL_TOP  = 2.45
LABEL_H    = 0.55

CIRCLE_D   = 1.1  # circle badge diameter
CARD_GAP   = 0.3
CARDS_TOP  = 3.1  # top of the number circle (card border starts CIRCLE_D/2 lower)
CARDS_H    = 6.4  # total height of cards region (including circle overhang)

HDR_TOP    = CARDS_TOP + CARDS_H + 0.4   # ~9.9
HDR_H      = 0.7
LINE_Y     = HDR_TOP + HDR_H             # ~10.6

ROWS_TOP   = LINE_Y + 0.3               # ~10.9
FOOTER_TOP = 18.2
ROWS_H     = FOOTER_TOP - ROWS_TOP      # ~7.3 cm available for use-case rows

# ── Table column widths ───────────────────────────────────────────────────────
C_UC  = 6.3   # Use Case
C_DI  = 3.8   # Data Inputs
C_MP  = 11.5  # Modelling Processes
C_DO  = CW - C_UC - C_DI - C_MP  # Data Outputs  (~9.67)


# ── Helpers ───────────────────────────────────────────────────────────────────

def _tf(shape):
    """Return text frame and set word wrap."""
    tf = shape.text_frame
    tf.word_wrap = True
    return tf


def add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(Cm(left), Cm(top), Cm(width), Cm(height))


def set_run(run, text, size_pt, bold=False, italic=False, color=BLACK):
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def add_para(tf, text, size_pt, bold=False, italic=False,
             color=BLACK, align=PP_ALIGN.LEFT, space_before=0):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    set_run(run, text, size_pt, bold=bold, italic=italic, color=color)
    return p


def add_shape(slide, shape_type, left, top, width, height):
    return slide.shapes.add_shape(shape_type, Cm(left), Cm(top), Cm(width), Cm(height))


def solid_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def no_fill(shape):
    shape.fill.background()


def no_line(shape):
    shape.line.fill.background()


def set_line(shape, color, width_pt):
    shape.line.color.rgb = color
    shape.line.width = Pt(width_pt)


def add_line(slide, x1, y1, x2, y2, color=BLACK, width_pt=1.0, dashed=False):
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT,
        Cm(x1), Cm(y1), Cm(x2), Cm(y2),
    )
    conn.line.color.rgb = color
    conn.line.width = Pt(width_pt)
    if dashed:
        _set_dashed(conn)
    return conn


def _set_dashed(connector):
    """Inject prstDash=dot into the connector's spPr/ln element."""
    try:
        spPr = connector._element.spPr
        ln = spPr.find(qn("a:ln"))
        if ln is None:
            ln = etree.SubElement(spPr, qn("a:ln"))
        for old in ln.findall(qn("a:prstDash")):
            ln.remove(old)
        prstDash = etree.SubElement(ln, qn("a:prstDash"))
        prstDash.set("val", "dot")
    except Exception:
        pass  # graceful degradation to solid line


# ── Slide builder ─────────────────────────────────────────────────────────────

def build_slide(prs, lang, dim, slide_num):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    prereqs   = dim["prerequisites"]
    use_cases = dim["useCases"]

    # prereq globalId → 1-based position number
    prereq_num = {p["globalId"]: i + 1 for i, p in enumerate(prereqs)}

    n_prereqs = len(prereqs)
    card_gap  = CARD_GAP
    card_w    = (CW - card_gap * (n_prereqs - 1)) / n_prereqs

    # ── Title ─────────────────────────────────────────────────────────────────
    tb = add_textbox(slide, ML, TITLE_TOP, CW, TITLE_H)
    tf = _tf(tb)
    p = tf.paragraphs[0]
    run = p.add_run()
    set_run(run, dim["name"][lang], size_pt=36, bold=True)

    # ── "Required Data Inputs" label ──────────────────────────────────────────
    section_label = {"en": "Required Data Inputs", "zh": "所需資料輸入"}[lang]
    tb = add_textbox(slide, ML, LABEL_TOP, CW, LABEL_H)
    tf = _tf(tb)
    p = tf.paragraphs[0]
    run = p.add_run()
    set_run(run, section_label, size_pt=11, bold=True, italic=True)

    # ── Prerequisite Cards ────────────────────────────────────────────────────
    card_border_top = CARDS_TOP + CIRCLE_D / 2      # card rect starts here
    card_border_h   = CARDS_H  - CIRCLE_D / 2       # card rect height

    for i, prereq in enumerate(prereqs):
        cx = ML + i * (card_w + card_gap)

        # Card border (rounded rectangle, white fill, black border)
        card = add_shape(
            slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            cx, card_border_top, card_w, card_border_h
        )
        solid_fill(card, WHITE)
        set_line(card, BLACK, 1.5)
        card.text_frame.word_wrap = True  # no text in the shape itself

        # Remove default text
        card.text_frame.paragraphs[0].clear()

        # Number circle badge
        circle_x = cx + card_w / 2 - CIRCLE_D / 2
        circle_y = CARDS_TOP
        circle = add_shape(
            slide, MSO_AUTO_SHAPE_TYPE.OVAL,
            circle_x, circle_y, CIRCLE_D, CIRCLE_D
        )
        solid_fill(circle, BLACK)
        no_line(circle)

        tf_c = circle.text_frame
        tf_c.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf_c.word_wrap = False
        p = tf_c.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        set_run(run, str(i + 1), size_pt=12, bold=True, color=WHITE)

        # Card content text box
        pad     = 0.3
        txt_top = card_border_top + CIRCLE_D / 2 + 0.15
        txt_h   = card_border_h   - CIRCLE_D / 2 - 0.3

        tb = add_textbox(slide, cx + pad, txt_top, card_w - pad * 2, txt_h)
        tf = _tf(tb)
        tf.vertical_anchor = MSO_ANCHOR.TOP

        # Card name
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        set_run(run, prereq["displayName"][lang], size_pt=10, bold=True)

        # Card items
        items = prereq.get("items", [])
        item_text = "\n".join(item[lang] for item in items)
        if item_text:
            add_para(tf, item_text, size_pt=9, align=PP_ALIGN.CENTER, space_before=5)

    # ── Column Headers ────────────────────────────────────────────────────────
    col_labels = {
        "en": ["Use Case", "Data Inputs", "Modelling Processes", "Data Outputs"],
        "zh": ["應用場景",  "資料輸入",     "建模流程",             "資料輸出"],
    }[lang]

    col_x      = [ML, ML + C_UC, ML + C_UC + C_DI, ML + C_UC + C_DI + C_MP]
    col_widths = [C_UC, C_DI, C_MP, C_DO]

    for i, (cx, cw, label) in enumerate(zip(col_x, col_widths, col_labels)):
        align = PP_ALIGN.LEFT if i == 0 else PP_ALIGN.CENTER
        tb = add_textbox(slide, cx, HDR_TOP, cw, HDR_H)
        tf = _tf(tb)
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        set_run(run, label, size_pt=11, bold=True)

    # Arrow triangles between headers
    for i in range(3):
        arrow_x = col_x[i] + col_widths[i] - 0.5
        tb = add_textbox(slide, arrow_x, HDR_TOP, 0.8, HDR_H)
        tf = _tf(tb)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        set_run(run, "▶", size_pt=9)

    # Horizontal divider
    add_line(slide, ML, LINE_Y, ML + CW, LINE_Y, color=BLACK, width_pt=1.5)

    # ── Use Case Rows ─────────────────────────────────────────────────────────
    n_rows = len(use_cases)
    row_h  = ROWS_H / n_rows

    for j, uc in enumerate(use_cases):
        row_top = ROWS_TOP + j * row_h

        # Dotted row separator (skip first)
        if j > 0:
            add_line(slide, ML, row_top, ML + CW, row_top,
                     color=LIGHT_GRAY, width_pt=0.75, dashed=True)

        pad_top = 0.25
        text_h  = row_h - 0.3

        def row_tb(col_idx, left, width, align=PP_ALIGN.LEFT):
            tb = add_textbox(slide, left, row_top + pad_top, width - 0.2, text_h)
            tf = _tf(tb)
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = align
            return tf

        # Use Case name
        tf = row_tb(0, col_x[0], C_UC)
        run = tf.paragraphs[0].add_run()
        set_run(run, uc["name"][lang], size_pt=10, bold=True)

        # Data Inputs (prereq numbers)
        req_ids = uc.get("requiredPrereqIds", [])
        nums    = sorted(prereq_num[rid] for rid in req_ids if rid in prereq_num)
        di_text = " + ".join(str(n) for n in nums)
        tf = row_tb(1, col_x[1], C_DI, align=PP_ALIGN.CENTER)
        run = tf.paragraphs[0].add_run()
        set_run(run, di_text, size_pt=12, bold=True)

        # Modelling Process
        tf = row_tb(2, col_x[2], C_MP, align=PP_ALIGN.CENTER)
        run = tf.paragraphs[0].add_run()
        set_run(run, uc["modellingProcess"][lang], size_pt=10)

        # Data Outputs
        tf = row_tb(3, col_x[3], C_DO, align=PP_ALIGN.CENTER)
        run = tf.paragraphs[0].add_run()
        set_run(run, uc["modelOutput"][lang], size_pt=10)

    # ── Footer ────────────────────────────────────────────────────────────────
    tb = add_textbox(slide, ML, FOOTER_TOP, 14.0, 0.7)
    tf = _tf(tb)
    p = tf.paragraphs[0]
    run = p.add_run()
    set_run(run, "FT STRATEGIES  |  Google News Initiative", size_pt=8, color=GRAY)

    conf_label = {"en": "Confidential", "zh": "專有與機密"}[lang]
    tb = add_textbox(slide, ML + CW - 6.0, FOOTER_TOP, 4.0, 0.7)
    tf = _tf(tb)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    set_run(run, conf_label, size_pt=8, color=GRAY)

    # Slide number
    tb = add_textbox(slide, ML + CW - 2.0, FOOTER_TOP, 2.0, 0.7)
    tf = _tf(tb)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    set_run(run, str(slide_num), size_pt=8, color=GRAY)


# ── Main ──────────────────────────────────────────────────────────────────────

def load_data():
    with open(DATA_DIR / "dimensions.json", encoding="utf-8") as f:
        return json.load(f)


def main():
    dims_data  = load_data()
    dimensions = dims_data["dimensions"]

    for lang, suffix in [("en", "EN"), ("zh", "ZH")]:
        prs = Presentation()
        prs.slide_width  = Cm(SLIDE_W)
        prs.slide_height = Cm(SLIDE_H)

        for i, dim in enumerate(dimensions, start=1):
            build_slide(prs, lang, dim, slide_num=i)

        out = Path(__file__).parent / f"AIL_Readiness_Assessment_{suffix}.pptx"
        prs.save(out)
        print(f"Saved: {out}")


if __name__ == "__main__":
    main()
